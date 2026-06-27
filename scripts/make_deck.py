#!/usr/bin/env python3
# Copyright (C) 2026 / Atsign migration
# SPDX-License-Identifier: Apache-2.0
"""
Generate slides/smart-route-atsign.pptx — the why/what + before/after security deck.

Run:  python scripts/make_deck.py
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Atsign-ish palette
INK = RGBColor(0x10, 0x18, 0x28)
TEAL = RGBColor(0x0E, 0x8C, 0x8C)
GREEN = RGBColor(0x13, 0xB5, 0x13)
RED = RGBColor(0xC0, 0x26, 0x26)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(p, text, size, color=INK, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_slide(title, subtitle, footer):
    s = prs.slides.add_slide(BLANK)
    bg(s, INK)
    tf = _box(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2))
    _run(tf.paragraphs[0], title, 40, WHITE, bold=True)
    p = tf.add_paragraph(); _run(p, subtitle, 22, RGBColor(0x9F, 0xE7, 0xE7))
    p2 = tf.add_paragraph(); p2.space_before = Pt(18); _run(p2, footer, 14, RGBColor(0xB8, 0xC0, 0xCC))
    # accent bar
    bar = s.shapes.add_shape(1, Inches(0.9), Inches(2.2), Inches(2.6), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()


def bullet_slide(title, bullets, accent=TEAL):
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    head = s.shapes.add_shape(1, 0, 0, W, Inches(1.15))
    head.fill.solid(); head.fill.fore_color.rgb = accent; head.line.fill.background()
    htf = head.text_frame; htf.margin_left = Inches(0.6)
    htf.paragraphs[0].alignment = PP_ALIGN.LEFT
    _run(htf.paragraphs[0], title, 26, WHITE, bold=True)

    tf = _box(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
    first = True
    for b in bullets:
        level = 0
        text = b
        if isinstance(b, tuple):
            level, text = b
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        p.level = level
        prefix = "•  " if level == 0 else "–  "
        _run(p, prefix + text, 20 if level == 0 else 17,
             INK if level == 0 else GREY, bold=(level == 0 and text.endswith(":")))


def compare_slide(title, rows):
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    head = s.shapes.add_shape(1, 0, 0, W, Inches(1.15))
    head.fill.solid(); head.fill.fore_color.rgb = INK; head.line.fill.background()
    _run(head.text_frame.paragraphs[0], title, 26, WHITE, bold=True)
    head.text_frame.margin_left = Inches(0.6)

    nrows = len(rows) + 1
    tbl_shape = s.shapes.add_table(nrows, 3, Inches(0.5), Inches(1.45),
                                   Inches(12.3), Inches(5.6))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.3)
    tbl.columns[1].width = Inches(4.5)
    tbl.columns[2].width = Inches(4.5)
    hdr = ["Aspect", "Before  (Intel demo as-is)", "After  (Atsign Platform)"]
    hcolors = [GREY, RED, GREEN]
    for c, (txt, col) in enumerate(zip(hdr, hcolors)):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = col
        tfc = cell.text_frame; tfc.word_wrap = True
        _run(tfc.paragraphs[0], txt, 14, WHITE, bold=True)
    for r, (aspect, before, after) in enumerate(rows, start=1):
        for c, txt in enumerate((aspect, before, after)):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r % 2 else WHITE
            tfc = cell.text_frame; tfc.word_wrap = True
            _run(tfc.paragraphs[0], txt, 12,
                 INK if c == 0 else (RED if c == 1 else GREEN),
                 bold=(c == 0))


def section_slide(title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    bg(s, TEAL)
    tf = _box(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.6))
    _run(tf.paragraphs[0], title, 34, WHITE, bold=True)
    if subtitle:
        p = tf.add_paragraph(); _run(p, subtitle, 18, RGBColor(0xE6, 0xFA, 0xFA))


# ---------------------------------------------------------------- slides
title_slide(
    "Smart Route Planning on the Atsign Platform",
    "Migrating Intel's edge-AI route planner to an identity-based, encrypted, port-less network",
    "Keep the AI brain · swap the nervous system · prove it end-to-end",
)

bullet_slide("The use case", [
    "Intel's Smart Route Planning Agent: an AI that routes a driver around congestion,",
    (1, "weather, road closures and accidents — in real time."),
    "A planning agent (LangGraph) consults distributed traffic-intersection agents and",
    (1, "data feeds, then picks and continuously re-optimises the route."),
    "Great on one laptop — but built in a way that cannot leave it safely.",
], accent=TEAL)

bullet_slide("The problem (as-is)", [
    "Every intersection must expose an inbound network port for the planner to poll.",
    "Connections are plain HTTP — no encryption, no authentication.",
    "Participants are found via a hand-edited list of localhost addresses — doesn't scale.",
    "No real driver experience; no notion of 'my route' per person.",
    "Weather / traffic-trends / planned-events optimisers ship turned OFF.",
], accent=RED)

bullet_slide("Security posture — BEFORE", [
    "Open inbound ports: one (or more) per intersection — a public attack surface.",
    "In the clear: unencrypted HTTP; data readable/modifiable in transit.",
    "No identity: nothing proves who a message is from; nothing to authorise against.",
    "IP/host-based: brittle, and exposed to scanning, spoofing, and NAT/firewall holes.",
    "Static trust: access = 'whatever is on the host list'.",
], accent=RED)

bullet_slide("What we built", [
    "Every participant gets its own secure identity (an atSign).",
    "Encrypted, port-less, push messaging between identities — devices only connect OUT.",
    "A zero-trust policy engine (default-deny) authorises every publisher/client by identity+role.",
    "Two front-ends: a Flutter commuter app (live reroute alerts) + a reused Gradio operator console.",
    "Intel's LangGraph reasoning reused UNCHANGED — only the transport was swapped.",
], accent=GREEN)

bullet_slide("Security posture — AFTER", [
    "Zero inbound ports: every device connects out to its own atServer — nothing to port-scan.",
    "End-to-end encrypted by default: only the addressed identity can read a message.",
    "Cryptographic identity: every message provably from a known atSign — no passwords, no tokens.",
    "Default-deny policy plane: access governed centrally by identity + role (segregation of duties).",
    "Dynamic + safe: a new intersection joins live once authorised — no host list, no restart.",
], accent=GREEN)

compare_slide("Before vs After — security posture", [
    ("Inbound ports", "Open port per intersection", "Zero — outbound-only connections"),
    ("Encryption", "Plain HTTP, in the clear", "End-to-end encrypted by default"),
    ("Identity / auth", "None (anonymous HTTP)", "Cryptographic atSign identity"),
    ("Addressing", "IP : port (host list)", "By identity (atSign)"),
    ("Access control", "Implicit (on the list)", "Default-deny policy engine"),
    ("Onboarding", "Edit a config file", "Policy-gated, live, no restart"),
    ("Attack surface", "Scan / spoof / MITM", "Nothing listening; nothing in clear"),
])

bullet_slide("How: keep the brain, swap the nervous system", [
    "KEEP (unchanged):",
    (1, "LangGraph agent, route logic, GPX parsing, map renderer, Gradio UI, data models."),
    "SWAP (4 small controllers):",
    (1, "read an encrypted subscription cache instead of polling URLs / CSVs — interface unchanged."),
    "ADD (the platform layer):",
    (1, "identity transport, policy engine, publishers, and the Flutter commuter app."),
    "The migration touched a tiny, well-defined seam — low risk, reusable pattern.",
], accent=TEAL)

bullet_slide("What we proved (verified live)", [
    "Python planner ⇄ Dart mobile app interoperate (encryption verified both ways).",
    "Intel's planner runs on pushed data — no polling, no host list.",
    "A genuine live reroute from an encrypted push (oakland → sanbruno).",
    "Default-deny enforced: an un-granted intersection is dropped until authorised.",
    "A brand-new intersection joins the running network live — no restart, no config edit.",
    "Reroute + alert reach the driver's phone and the operator's map in real time.",
], accent=GREEN)

bullet_slide("Why it was worth it", [
    "Removes an entire class of attacks — no open ports, nothing in clear.",
    "Turns a single-laptop demo into a city-scale, governable architecture.",
    "Trust becomes first-class: central, role-based, segregation of duties.",
    "Re-activated capabilities Intel left disabled (weather / events as live feeds).",
    "Proved a low-risk, reusable template for moving other edge-AI suites onto the platform.",
    "Tangible: a phone that reroutes live and a control-room map — a security story you can see.",
], accent=TEAL)

bullet_slide("Status & next steps", [
    "Status: end-to-end working & demonstrable on the platform (Phases 0–9 high-value slice).",
    "One command brings the live stack up; a pure-Dart tool drives reroutes.",
    "Next (pilot): run the Flutter app on a device; extend policy to requests/viewing;",
    (1, "idempotency / reconnect hardening; point feeds at real data sources (no planner change)."),
], accent=INK if False else TEAL)

section_slide("Keep the AI. Change the network.",
              "Identity-based · encrypted · port-less · policy-governed — proven end-to-end.")

out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slides")
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "smart-route-atsign.pptx")
prs.save(out)
print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")
