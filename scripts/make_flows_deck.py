#!/usr/bin/env python3
# Copyright (C) 2026 / Atsign migration
# SPDX-License-Identifier: Apache-2.0
"""
Generate slides/smart-route-flows-vpn-costs.pptx — drawn architecture/flow diagrams:
separate edge machines doing camera inference publishing over the Atsign Platform,
the same system built the traditional VPN/firewall way, and a 3-year cost comparison.

Run:  python scripts/make_flows_deck.py
"""
import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# Atsign-ish palette (matches make_deck.py)
INK = RGBColor(0x10, 0x18, 0x28)
TEAL = RGBColor(0x0E, 0x8C, 0x8C)
TEAL_DK = RGBColor(0x0A, 0x6E, 0x6E)
GREEN = RGBColor(0x13, 0xB5, 0x13)
RED = RGBColor(0xC0, 0x26, 0x26)
RED_DK = RGBColor(0x8E, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)
MINT = RGBColor(0xE9, 0xF9, 0xEF)
ROSE = RGBColor(0xFB, 0xEE, 0xEE)
SKY = RGBColor(0xE6, 0xF4, 0xF4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


# ------------------------------------------------------------------ helpers
def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def header(slide, title, accent=TEAL, sub=""):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = accent; band.line.fill.background()
    tf = band.text_frame; tf.margin_left = Inches(0.55); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
    if sub:
        r2 = tf.paragraphs[0].add_run(); r2.text = "   " + sub
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xDE, 0xF2, 0xF2)


def card(slide, x, y, w, h, title, lines, fill=WHITE, edge=TEAL, title_color=None,
         title_size=12.5, line_size=10.5, line_color=GREY, bold_title=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = edge; shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.09); tf.margin_right = Inches(0.09)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(title_size); r.font.bold = bold_title
    r.font.color.rgb = title_color or edge
    for ln in lines:
        p = tf.add_paragraph(); p.space_before = Pt(1)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(line_size); r.font.color.rgb = line_color
    return shp


def band(slide, x, y, w, h, title, lines, fill, edge, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.05
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = edge; shp.line.width = Pt(1.75)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = text_color
    for ln in lines:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_before = Pt(4)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(10); r.font.color.rgb = text_color
    return shp


def arrow(slide, x1, y1, x2, y2, color=TEAL, width=2.25, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        dash = etree.SubElement(ln, qn("a:prstDash"))
        dash.set("val", "dash")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    conn.shadow.inherit = False
    return conn


def label(slide, x, y, w, text, size=10, color=GREY, bold=False, align=PP_ALIGN.LEFT,
          italic=False):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic
    return tb


def badge(slide, x, y, w, text, fill=MINT, edge=GREEN, color=INK, size=10.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.42))
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = edge; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    return shp


# ------------------------------------------------------------- slide 1: title
s = prs.slides.add_slide(BLANK)
bg(s, INK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.6), Inches(3.0))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Edge AI, City-Wide — Without a VPN"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph()
r = p.add_run()
r.text = ("Distributed camera-inference machines publishing over the Atsign Platform — "
          "the flows, the traditional VPN/firewall alternative, and the economics")
r.font.size = Pt(20); r.font.color.rgb = RGBColor(0x9F, 0xE7, 0xE7)
p = tf.add_paragraph(); p.space_before = Pt(18)
r = p.add_run()
r.text = "Smart Route Planning migration · verified end-to-end · Python + Dart SDKs"
r.font.size = Pt(14); r.font.color.rgb = RGBColor(0xB8, 0xC0, 0xCC)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.1), Inches(2.6), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()

# --------------------------------------- slide 2: deployment (the Atsign way)
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header(s, "Deployment — every machine is an identity, no machine listens",
       accent=TEAL, sub="each intersection runs its own inference box; only tiny encrypted messages leave the site")

EDGE_X, EDGE_W, EDGE_H = Inches(0.35), Inches(3.25), Inches(1.42)
ys = [Inches(1.30), Inches(2.86), Inches(4.42)]
names = [("Intersection — Market St", "@intc_intxn_market_st"),
         ("Intersection — 5th Ave", "@intc_intxn_5th_ave"),
         ("Intersection — Broadway  (…×N)", "@intc_intxn_broadway")]
for (title, at), y in zip(names, ys):
    card(s, EDGE_X, y, EDGE_W, EDGE_H, title,
         ["📷 camera → LOCAL inference (vehicle density)",
          "edge machine · Intel CPU/GPU · video never leaves site",
          f"publishes as  {at}"],
         fill=SKY, edge=TEAL_DK, title_size=12, line_size=9.5, line_color=INK)
card(s, EDGE_X, Inches(5.98), EDGE_W, Inches(0.98), "Condition feed machines",
     ["weather / traffic-trends / planned events",
      "@intc_weather_feed · @intc_traffic_trends_feed · @intc_events_feed"],
     fill=SKY, edge=TEAL_DK, title_size=12, line_size=9.5, line_color=INK)

PLAT_X, PLAT_W = Inches(4.28), Inches(1.62)
band(s, PLAT_X, Inches(1.30), PLAT_W, Inches(5.50), "Atsign Platform",
     ["atServers (cloud)", "end-to-end encrypted", "store & forward",
      "outbound-only TLS on BOTH sides", "no inbound ports anywhere",
      "addressed by identity, not IP"],
     fill=TEAL, edge=TEAL_DK)

PLN_X, PLN_W = Inches(6.62), Inches(3.05)
card(s, PLN_X, Inches(1.55), PLN_W, Inches(1.75), "Route Planner (server)",
     ["@intc_smartroute_planner",
      "Intel LangGraph brain — unchanged",
      "re-plans every ~8 s from the encrypted cache",
      "accepts ONLY policy-granted publishers"],
     fill=MINT, edge=GREEN, line_color=INK, title_size=13)
card(s, PLN_X, Inches(4.10), PLN_W, Inches(1.55), "Policy plane (default-deny)",
     ["engine  @intc_route_policy",
      "admin (web)  @intc_route_policy_admin",
      "grant / revoke a machine in one click — applies in seconds"],
     fill=LIGHT, edge=GREY, line_color=INK, title_size=13)

CON_X, CON_W = Inches(10.35), Inches(2.62)
card(s, CON_X, Inches(1.55), CON_W, Inches(1.45), "Operator console (web)",
     ["@intc_route_operator", "live map · 🚨 reroute status"],
     fill=MINT, edge=GREEN, line_color=INK, title_size=12.5)
card(s, CON_X, Inches(3.60), CON_W, Inches(1.55), "Commuter app (Flutter)",
     ["@intc_commuter01", "route + reroute alerts on the phone", "works on any network — no VPN client"],
     fill=MINT, edge=GREEN, line_color=INK, title_size=12.5)

# arrows: producers -> platform
for y in ys:
    arrow(s, EDGE_X + EDGE_W, y + Emu(int(EDGE_H) // 2), PLAT_X, y + Emu(int(EDGE_H) // 2), color=TEAL)
arrow(s, EDGE_X + EDGE_W, Inches(6.57), PLAT_X, Inches(6.57), color=TEAL)
# platform -> planner (ingest) and planner -> consumers (push, via platform)
arrow(s, PLAT_X + PLAT_W, Inches(2.35), PLN_X, Inches(2.35), color=TEAL)
arrow(s, PLN_X + PLN_W, Inches(2.10), CON_X, Inches(2.10), color=GREEN, dashed=True)
arrow(s, PLN_X + PLN_W, Inches(2.60), CON_X, Inches(4.20), color=GREEN, dashed=True)
label(s, Inches(9.72), Inches(3.05), Inches(1.7), "route + status pushes (via platform)", size=9,
      color=GREEN, bold=True)
arrow(s, PLN_X + Inches(1.5), Inches(4.10), PLN_X + Inches(1.5), Inches(3.30), color=GREY, width=1.75)
label(s, Inches(8.25), Inches(3.62), Inches(1.6), "grants", size=9, color=GREY, bold=True)

badge(s, Inches(0.35), Inches(7.02), Inches(2.25), "0 inbound ports", fill=MINT, edge=GREEN)
badge(s, Inches(2.75), Inches(7.02), Inches(2.6), "E2E encrypted by default", fill=MINT, edge=GREEN)
badge(s, Inches(6.62), Inches(7.02), Inches(3.5), "~1 KB messages — no video backhaul", fill=MINT, edge=GREEN)
badge(s, Inches(10.35), Inches(7.02), Inches(2.62), "join = a policy grant", fill=MINT, edge=GREEN)

# ------------------------------------------------- slide 3: the message flow
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header(s, "The flow — camera frame to rerouted commuter, end to end encrypted",
       accent=TEAL, sub="numbered path of one congestion event")

lane_y = [Inches(1.45), Inches(3.55), Inches(5.55)]
label(s, Inches(0.35), lane_y[0] - Inches(0.02), Inches(4.0), "AT THE INTERSECTION (edge machine)",
      size=11, color=TEAL_DK, bold=True)
label(s, Inches(0.35), lane_y[1] - Inches(0.02), Inches(4.0), "AT THE PLANNER (server)",
      size=11, color=TEAL_DK, bold=True)
label(s, Inches(0.35), lane_y[2] - Inches(0.02), Inches(4.0), "AT THE CONSUMERS",
      size=11, color=TEAL_DK, bold=True)

STEP_H = Inches(1.28)
steps_top = [
    ("1 · Capture", ["camera frames analysed on-site", "video never leaves the machine"]),
    ("2 · Infer", ["vehicle density = 30 (threshold 10)", "OpenVINO/CPU — no cloud round-trip"]),
    ("3 · Publish", ["notify: live_traffic → planner", "encrypted to the planner's key,",
                     "signed by the intersection's atSign", "TTL 60 s — stale data self-erases"]),
]
xs = [Inches(0.50), Inches(4.85), Inches(9.20)]
w3 = Inches(4.00)
for (t, lines), x in zip(steps_top, xs):
    card(s, x, lane_y[0] + Inches(0.32), w3, STEP_H, t, lines,
         fill=SKY, edge=TEAL_DK, line_color=INK, title_size=13, line_size=10)
arrow(s, xs[0] + w3, lane_y[0] + Inches(0.95), xs[1], lane_y[0] + Inches(0.95))
arrow(s, xs[1] + w3, lane_y[0] + Inches(0.95), xs[2], lane_y[0] + Inches(0.95))

steps_mid = [
    ("4 · Verify & admit", ["provably from @intc_intxn_market_st", "policy check: granted? (default-deny)",
                            "revoked publisher → dropped + cache purged"]),
    ("5 · Re-plan", ["Intel LangGraph runs unchanged", "density > threshold on current route",
                     "→ pick next shortest clean route"]),
    ("6 · Push", ["route → commuter, status → operator", "encrypted per recipient",
                  "no consumer polls anything"]),
]
for (t, lines), x in zip(steps_mid, xs):
    card(s, x, lane_y[1] + Inches(0.32), w3, STEP_H, t, lines,
         fill=MINT, edge=GREEN, line_color=INK, title_size=13, line_size=10)
arrow(s, xs[0] + w3, lane_y[1] + Inches(0.95), xs[1], lane_y[1] + Inches(0.95), color=GREEN)
arrow(s, xs[1] + w3, lane_y[1] + Inches(0.95), xs[2], lane_y[1] + Inches(0.95), color=GREEN)
# lane hop arrows
arrow(s, xs[2] + Inches(1.6), lane_y[0] + Inches(0.32) + STEP_H, xs[0] + Inches(0.4),
      lane_y[1] + Inches(0.32), color=TEAL, width=2.0, dashed=True)
arrow(s, xs[2] + Inches(1.6), lane_y[1] + Inches(0.32) + STEP_H, xs[0] + Inches(0.4),
      lane_y[2] + Inches(0.32), color=GREEN, width=2.0, dashed=True)

steps_bot = [
    ("7 · Operator console", ["map flips to 🚨 REROUTE", "watchdog-hardened subscriber"]),
    ("8 · Commuter phone", ["new route + alert appear live", "no VPN app, any network"]),
    ("9 · Self-clearing", ["TTL expires or density drops", "→ route reverts, alert clears"]),
]
for (t, lines), x in zip(steps_bot, xs):
    card(s, x, lane_y[2] + Inches(0.32), w3, Inches(1.05), t, lines,
         fill=LIGHT, edge=GREY, line_color=INK, title_size=13, line_size=10)

label(s, Inches(0.35), Inches(7.05), Inches(12.6),
      "Every hop is identity-to-identity and encrypted; a compromised network in the middle learns nothing and can inject nothing.",
      size=11.5, color=INK, bold=True, align=PP_ALIGN.CENTER)

# ------------------------------------ slide 4: the traditional VPN/firewall way
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header(s, "The traditional way — VPN concentrator, firewalls, and N listening endpoints",
       accent=RED, sub="same system, built the way it's usually built")

for (title, _), y in zip(names, ys):
    card(s, EDGE_X, y, EDGE_W, EDGE_H, title.replace("  (…×N)", " (…×N)"),
         ["📷 camera + inference box",
          "static IP (or DDNS) · IPsec firewall/router",
          "OPEN endpoint: IKE/500, ESP, or TLS/443 listening"],
         fill=ROSE, edge=RED, line_color=INK, title_size=12, line_size=9.5)
card(s, EDGE_X, Inches(5.98), EDGE_W, Inches(0.98), "Feed sources",
     ["per-source tunnels or an exposed API gateway", "more certs, more rules, more tickets"],
     fill=ROSE, edge=RED, line_size=9.5, line_color=INK)

band(s, Inches(4.28), Inches(1.30), Inches(1.62), Inches(5.50), "Internet",
     ["N site-to-site IPsec tunnels", "PSK / cert per site", "NAT traversal pain",
      "every endpoint scannable"],
     fill=GREY, edge=INK)

card(s, Inches(6.62), Inches(1.55), Inches(3.05), Inches(2.00), "HQ — corporate firewall",
     ["VPN concentrator (HA pair)",
      "single point of failure AND a choke point",
      "appliance CVEs → emergency patch fire-drills",
      "every packet hairpins through here"],
     fill=ROSE, edge=RED_DK, line_color=INK, title_size=13)
card(s, Inches(6.62), Inches(4.10), Inches(3.05), Inches(1.30), "Route planner (behind FW)",
     ["reachable only inside the VPN",
      "consumers must come to IT — or IT opens a DMZ"],
     fill=LIGHT, edge=GREY, line_color=INK, title_size=12.5)

card(s, CON_X, Inches(1.55), CON_W, Inches(1.45), "Operator console",
     ["needs VPN access / on-net PC"],
     fill=ROSE, edge=RED, line_color=INK, title_size=12.5)
card(s, CON_X, Inches(3.60), CON_W, Inches(1.55), "Commuter app",
     ["per-device VPN client licences", "…or a public DMZ API to secure", "on-boarding friction per user"],
     fill=ROSE, edge=RED, line_color=INK, title_size=12.5)

for y in ys:
    arrow(s, EDGE_X + EDGE_W, y + Emu(int(EDGE_H) // 2), Inches(4.28), y + Emu(int(EDGE_H) // 2), color=RED)
arrow(s, EDGE_X + EDGE_W, Inches(6.57), Inches(4.28), Inches(6.57), color=RED)
arrow(s, Inches(5.90), Inches(2.55), Inches(6.62), Inches(2.55), color=RED)
arrow(s, Inches(8.14), Inches(3.55), Inches(8.14), Inches(4.10), color=RED_DK, width=1.75)
arrow(s, Inches(9.67), Inches(2.30), CON_X, Inches(2.30), color=RED, dashed=True)
arrow(s, Inches(9.67), Inches(2.80), CON_X, Inches(4.35), color=RED, dashed=True)

badge(s, Inches(0.35), Inches(7.02), Inches(3.1), "N+1 listening attack surfaces", fill=ROSE, edge=RED)
badge(s, Inches(3.60), Inches(7.02), Inches(3.1), "change ticket per site, per change", fill=ROSE, edge=RED)
badge(s, Inches(6.85), Inches(7.02), Inches(2.9), "concentrator = SPOF", fill=ROSE, edge=RED)
badge(s, Inches(9.90), Inches(7.02), Inches(3.05), "cert/PSK rotation ×N sites", fill=ROSE, edge=RED)

# -------------------------------------------------- slide 5: comparison table
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header(s, "Side by side — the same system, two networks", accent=INK)
rows = [
    ("Connectivity model", "Hub-and-spoke IPsec tunnels via concentrator", "Direct identity↔identity, outbound-only both ends"),
    ("What listens at a site", "IKE/ESP or TLS endpoint on every site", "Nothing — zero inbound ports to scan"),
    ("Addressing", "Static IPs / DDNS + routing tables", "By identity (@atSign); IPs irrelevant, cellular fine"),
    ("Encryption", "Encrypted in the tunnel, plaintext beyond it", "End-to-end: only the addressed identity can decrypt"),
    ("Add an intersection", "Ticket → firewall change → tunnel + certs → route update", "Activate an atSign → one policy grant → live in minutes"),
    ("Revoke a machine", "Track down rules/certs across boxes", "One click in the policy console; enforced in seconds"),
    ("Mobile / 3rd-party users", "VPN client licences or a DMZ API to build & defend", "Just the app — encrypted pushes reach any network"),
    ("Failure modes", "Concentrator outage = city-wide outage", "No central data-path chokepoint to lose"),
    ("Patching burden", "Appliance CVE fire-drills (VPN CVEs are attackers' #1 door)", "No network appliances to patch"),
]
tbl_shape = s.shapes.add_table(len(rows) + 1, 3, Inches(0.4), Inches(1.30),
                               Inches(12.5), Inches(5.8))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.7)
tbl.columns[1].width = Inches(5.0)
tbl.columns[2].width = Inches(4.8)
for c, (txt, col) in enumerate(zip(["Aspect", "Traditional VPN / firewall", "Atsign Platform"],
                                   [GREY, RED, GREEN])):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = col
    r = cell.text_frame.paragraphs[0].add_run(); r.text = txt
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
for ri, (a, b, c_) in enumerate(rows, start=1):
    for ci, txt in enumerate((a, b, c_)):
        cell = tbl.cell(ri, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if ri % 2 else WHITE
        cell.text_frame.word_wrap = True
        r = cell.text_frame.paragraphs[0].add_run(); r.text = txt
        r.font.size = Pt(11)
        r.font.bold = (ci == 0)
        r.font.color.rgb = INK if ci == 0 else (RED_DK if ci == 1 else RGBColor(0x0B, 0x7A, 0x0B))

# ------------------------------------------------------- slide 6: cost model
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header(s, "What it costs — indicative 3-year model, 25 intersections + 3 feeds + planner",
       accent=INK, sub="the structure is the point; substitute your own quotes")
cost_rows = [
    ("VPN concentrator (HA pair) + support", "$30,000", "$0  — none needed"),
    ("Per-site IPsec firewall/router + support (25×)", "$24,000", "$0  — any outbound internet (cellular OK)"),
    ("Static IPs / business-line premium (28×)", "$5,000", "$0  — no inbound reachability required"),
    ("Firewall & VPN operations (~0.25 FTE: tickets, tunnels, certs)", "$90,000", "$18,000  (~0.05 FTE: policy grants in a console)"),
    ("Appliance patching, CVE fire-drills, site visits", "$15,000", "$0  — no network appliances"),
    ("Platform subscription (≈30 identities, indicative)", "$0", "$11,000"),
    ("TOTAL (3 years)", "$164,000", "$29,000"),
    ("Per intersection, per year", "≈ $2,190", "≈ $385"),
]
tbl_shape = s.shapes.add_table(len(cost_rows) + 1, 3, Inches(0.4), Inches(1.35),
                               Inches(12.5), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(5.6)
tbl.columns[1].width = Inches(3.45)
tbl.columns[2].width = Inches(3.45)
for c, (txt, col) in enumerate(zip(["Cost line (3-yr)", "Traditional VPN / firewall", "Atsign Platform"],
                                   [GREY, RED, GREEN])):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = col
    r = cell.text_frame.paragraphs[0].add_run(); r.text = txt
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
for ri, (a, b, c_) in enumerate(cost_rows, start=1):
    total = a.startswith("TOTAL") or a.startswith("Per intersection")
    for ci, txt in enumerate((a, b, c_)):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = (RGBColor(0xE8, 0xEC, 0xF2) if total
                                    else (LIGHT if ri % 2 else WHITE))
        cell.text_frame.word_wrap = True
        r = cell.text_frame.paragraphs[0].add_run(); r.text = txt
        r.font.size = Pt(12 if not total else 13)
        r.font.bold = (ci == 0) or total
        r.font.color.rgb = INK if ci == 0 else (RED_DK if ci == 1 else RGBColor(0x0B, 0x7A, 0x0B))
badge(s, Inches(3.4), Inches(6.30), Inches(6.5),
      "≈ 5–6× lower 3-year TCO — and the gap WIDENS with every site added",
      fill=MINT, edge=GREEN, size=13)
label(s, Inches(0.4), Inches(6.90), Inches(12.5),
      "Not counted (all favour the edge/atSign model): no video backhaul bandwidth, no per-user VPN licences, "
      "no DMZ API build-out, avoided breach exposure from N listening VPN endpoints.",
      size=10.5, color=GREY, italic=True)

# ------------------------------------------------ slide 7: advantages summary
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
header(s, "Why this wins — beyond the invoice", accent=TEAL)
adv = [
    ("Security you can't misconfigure", "No listening endpoints to scan; E2E encryption is the default, not a firewall rule that can drift."),
    ("Scale is a grant, not a project", "Site #26 = activate an identity + one policy click. No tunnel, no ticket, no truck roll."),
    ("Edge inference stays at the edge", "Video is processed on the intersection machine; only ~1 KB encrypted facts travel. 4G/5G is enough."),
    ("No chokepoint, no SPOF", "Messages flow identity-to-identity; there is no concentrator to size, patch, or lose."),
    ("Governance is first-class", "Default-deny policy plane, one console: who may publish, provable per message, revocable in seconds."),
    ("The brain was untouched", "Intel's LangGraph planner ran unchanged — this is a transport swap, a repeatable pattern for any edge-AI suite."),
]
y = Inches(1.35)
for i, (t, d) in enumerate(adv):
    x = Inches(0.5) if i % 2 == 0 else Inches(6.85)
    card(s, x, y, Inches(6.0), Inches(1.62), t, [d],
         fill=MINT if i % 2 == 0 else SKY, edge=GREEN if i % 2 == 0 else TEAL_DK,
         line_color=INK, title_size=14, line_size=11.5)
    if i % 2 == 1:
        y += Inches(1.86)

# ------------------------------------------------------------ slide 8: close
s = prs.slides.add_slide(BLANK)
bg(s, TEAL)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Put the intelligence at the edge. Put the trust in identity."
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph()
r = p.add_run()
r.text = ("Camera inference on independent machines · encrypted identity-addressed messages · "
          "no VPN, no firewall rules, no inbound ports — at a fraction of the cost.")
r.font.size = Pt(17); r.font.color.rgb = RGBColor(0xE6, 0xFA, 0xFA)

# --------------------------------------------------------------------- save
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slides")
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "smart-route-flows-vpn-costs.pptx")
prs.save(out)
print(f"wrote {out} ({len(prs.slides._sldIdLst)} slides)")
